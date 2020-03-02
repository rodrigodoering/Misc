import pandas as pd
import os
import string
import nltk
import re
import pandas as pd
from nltk.tokenize import word_tokenize
from textdistance import jaro_winkler as jk
from pptx import Presentation
from SQLServer import SQLServer

vocab = pd.read_excel('vocabulário_final.xlsx')
verbos = pd.read_csv('tb_verbos.csv')
stopwords = nltk.corpus.stopwords.words("portuguese")

class ppt(object):

    def __init__(self, path):
        self.file = Presentation(path)
        self.texto_processado = False 
    
    @staticmethod
    def header(_string_):
        ruido = ['mba', 'aula', 'live', 'pós', 'graduação', 'capacitação', 'workshop', 'especialização', 'atividade', 'http', 'www']
        return bool(sum([1 if termo in _string_.lower() else 0 for termo in ruido]))

    @staticmethod                                                           
    def padronizar_texto(array_strings):
        texto = list(map(lambda x: x.lower().strip(string.punctuation).replace(",", ""), array_strings))
        return [word_tokenize(sentenca, language="portuguese") for sentenca in texto if len(sentenca) > 0]

    @staticmethod
    def split_terms(array_termos):
        novos_termos = []
        remover = []
        
        for termo in array_termos:
            if '/' in termo:
                remover.append(termo)
                splitted = termo.split('/')
                
                for termo_novo in splitted:
                    novos_termos.append(termo_novo)
                
        array_termos = [termo for termo in array_termos if termo not in remover]
        
        for termo_novo in novos_termos:
            array_termos.append(termo_novo)
        
        return [termo for termo in array_termos if termo not in stopwords and len(termo) > 1]


    def processar_texto(self):
        texto = []

        for slide in self.file.slides:

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if not ppt.header(shape.text):
                        texto.append(shape.text)
        
        texto_filtrado = [termo for array_termos in ppt.padronizar_texto(texto) for termo in array_termos if termo not in stopwords and len(termo) > 1]
        vocabulario_final = []

        for termo in texto_filtrado:
            inicio, fim, carac = termo[:2], termo[-2:], len(termo)
            df_verbos = verbos[(verbos.inicio == inicio) & (verbos.caracteres == carac) & (verbos.fim == fim)]
            
            if df_verbos.shape[0] == 0:
                vocabulario_final.append(termo)

            elif termo not in df_verbos.verbo.values.tolist():
                vocabulario_final.append(termo)
            
            else:
                pass
        
        self.texto_processado = list(set(ppt.split_terms(vocabulario_final)))


    def mapear(self):
        if not self.texto_processado:
            print('Necessário processar o texto antes de mapear')
            return
        
        else:
            representacao_final = []
            
            for termo in vocab.termo:
                distancias = [jk.distance(termo, termo_ppt) for termo_ppt in self.texto_processado]
                busted = [termo_ppt for termo_ppt, dist in zip(texto, distancias) if dist <= .1]
                
                if len(busted) == 0:
                    pass
                
                else:
                    for termo_ in busted:
            
                        representacao_final.append(termo)
            
            representacao = list(set(representacao_final))
            return ','.join(representacao)

